#!/usr/bin/env python3
"""
Murdoch University — AI-Powered Prospective Student Experience
Pitch Deck Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Palette (from Tailwind config) ──────────────────────────────────────
RED         = RGBColor(0xC8, 0x20, 0x3A)   # primary-600 — Murdoch Red
DARK_RED    = RGBColor(0x8C, 0x0E, 0x26)   # primary-800
PINK        = RGBColor(0xEE, 0x55, 0x90)   # pink-600
DARK_BG     = RGBColor(0x21, 0x21, 0x21)   # grey-900
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
AQUA        = RGBColor(0x00, 0x9C, 0x98)   # aqua-600
AQUA_LIGHT  = RGBColor(0xD4, 0xF4, 0xF3)   # aqua-100
LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)   # grey-100
GREY_300    = RGBColor(0xE0, 0xE0, 0xE0)   # grey-300
MID_GREY    = RGBColor(0x9E, 0x9E, 0x9E)   # grey-500
DARK_GREY   = RGBColor(0x61, 0x61, 0x61)   # grey-700
NEAR_BLACK  = RGBColor(0x21, 0x21, 0x21)   # grey-900
PURPLE      = RGBColor(0x4D, 0x32, 0x99)   # purple-800
PURPLE_LIGHT= RGBColor(0xDD, 0xD4, 0xF5)   # purple-200
RED_LIGHT   = RGBColor(0xF5, 0xC3, 0xCB)   # primary-200
ORANGE      = RGBColor(0xFF, 0x84, 0x00)   # orange-600
ORANGE_LIGHT= RGBColor(0xFF, 0xD4, 0xA8)   # orange-200

# ── Dimensions (16:9) ─────────────────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.65)   # left margin
CW = SW - ML * 2    # content width

# ── Core Helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = wrap
    return box.text_frame

def p0(tf, text, size=14, bold=False, color=NEAR_BLACK,
       align=PP_ALIGN.LEFT, before=0, after=0, italic=False):
    """Write to the first (existing) paragraph of a fresh textbox."""
    p = tf.paragraphs[0]
    p.alignment = align
    if before: p.space_before = Pt(before)
    if after:  p.space_after  = Pt(after)
    r = p.add_run()
    r.text = text
    r.font.name  = "Poppins"
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color

def pn(tf, text, size=14, bold=False, color=NEAR_BLACK,
       align=PP_ALIGN.LEFT, before=0, after=0, italic=False):
    """Add a new paragraph to a textbox."""
    p = tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    if after:  p.space_after  = Pt(after)
    r = p.add_run()
    r.text = text
    r.font.name  = "Poppins"
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color

def title_bar(slide, title, subtitle=None):
    """Standard slide header: title + thin red rule."""
    f = tb(slide, ML, Inches(0.5), CW, Inches(0.7), wrap=False)
    p0(f, title, size=24, bold=True, color=NEAR_BLACK)
    if subtitle:
        pn(f, subtitle, size=13, color=DARK_GREY, before=2)
    add_rect(slide, ML, Inches(1.3), CW, Inches(0.05), RED)

def screenie(slide, l, t, w, h, label):
    add_rect(slide, l, t, w, h, LIGHT_GREY, GREY_300)
    f = tb(slide, l, t + h/2 - Inches(0.35), w, Inches(0.7))
    p0(f, f"[ {label} ]", size=10, color=MID_GREY,
       italic=True, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, bg=LIGHT_GREY, accent=None):
    r = add_rect(slide, l, t, w, h, bg)
    if accent:
        add_rect(slide, l, t, w, Inches(0.06), accent)
    return r

# ── Slide Builders ────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = blank(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, Inches(0.16), SH, RED)
    add_rect(sl, 0, 0, SW, Inches(0.08), RED)

    f = tb(sl, Inches(0.5), Inches(1.5), Inches(10), Inches(0.45))
    p0(f, "MURDOCH UNIVERSITY", size=11, bold=True, color=RED)

    f = tb(sl, Inches(0.5), Inches(2.2), Inches(11), Inches(2.2))
    p0(f, "AI-Powered Prospective", size=46, bold=True, color=WHITE)
    pn(f, "Student Experience", size=46, bold=True, color=WHITE)

    f = tb(sl, Inches(0.5), Inches(4.55), Inches(10), Inches(0.55))
    p0(f, "Turning website friction into guided conversations", size=20, color=MID_GREY)

    add_rect(sl, Inches(0.5), Inches(5.35), Inches(4.5), Inches(0.04), RED)

    f = tb(sl, Inches(0.5), Inches(5.6), Inches(9), Inches(0.4))
    p0(f, "Prepared for Murdoch University  —  March 2026", size=12, color=MID_GREY)


def slide_agenda(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Today's Agenda")

    items = [
        ("01", "The Challenge",        "What the research tells us about the prospective student journey"),
        ("02", "The Solution",         "A working AI agent built for Murdoch, ready to demo"),
        ("03", "From PoC to Production","Federating knowledge and building a platform"),
        ("04", "The Strategic Case",   "Why this matters commercially and architecturally"),
        ("05", "Why Now",              "The window for first-mover advantage"),
        ("06", "The Roadmap",          "A phased engagement from pilot to scale"),
    ]

    row_h = Inches(0.8)
    top   = Inches(1.5)

    for i, (num, heading, sub) in enumerate(items):
        t = top + i * row_h
        add_rect(sl, ML, t + Inches(0.1), Inches(0.42), Inches(0.42), RED)
        f = tb(sl, ML, t + Inches(0.08), Inches(0.42), Inches(0.42))
        p0(f, num, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        f = tb(sl, ML + Inches(0.56), t, Inches(11), Inches(0.38), wrap=False)
        p0(f, heading, size=16, bold=True, color=NEAR_BLACK)

        f = tb(sl, ML + Inches(0.56), t + Inches(0.37), Inches(11), Inches(0.35), wrap=False)
        p0(f, sub, size=12, color=DARK_GREY)

        if i < len(items) - 1:
            add_rect(sl, ML + Inches(0.56), t + row_h - Inches(0.06),
                     Inches(11.5), Inches(0.01), GREY_300)


def slide_section(prs, number, title, bg=DARK_BG, accent=RED):
    sl = blank(prs)
    fill_bg(sl, bg)

    # Faded large number in background
    f = tb(sl, SW - Inches(6), SH/2 - Inches(3), Inches(6), Inches(6))
    p0(f, number, size=200, bold=True,
       color=RGBColor(0x2D, 0x2D, 0x2D), align=PP_ALIGN.RIGHT)

    add_rect(sl, ML, SH/2 - Inches(0.55), Inches(1.6), Inches(0.06), accent)

    f = tb(sl, ML, SH/2 - Inches(0.35), Inches(8), Inches(0.45))
    p0(f, f"{number} /", size=14, bold=True, color=accent)

    f = tb(sl, ML, SH/2 + Inches(0.15), Inches(9), Inches(1.3))
    p0(f, title, size=52, bold=True, color=WHITE)


def slide_challenge(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Your Website Is Working Against Your Enrolment Funnel")

    points = [
        ("The split students feel but can't see",
         "The site is divided between two sub-sites with different owners — students bear the cost of that division."),
        ("Study sub-site (Marketing-owned)",
         "Conversion-focused and brand-consistent, but covers only part of the picture."),
        ("Explore sub-site (School/Faculty-owned)",
         "Deep course content — but every school maintains it differently. No unified voice or student-facing logic."),
        ("Invisible seams, real confusion",
         "Students cross between these environments without knowing it, hitting inconsistent layouts and terminology."),
        ("The structural result",
         "Navigation that maps to Murdoch's internal ownership structure, not how a prospective student thinks."),
        ("Every unanswered question",
         "Erodes confidence in Murdoch. That erosion is structural, not incidental."),
    ]

    col_w = Inches(5.6)
    gap   = Inches(0.4)

    for i, (heading, body) in enumerate(points):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        left = ML if col == 0 else ML + col_w + gap
        top  = Inches(1.5) + row * Inches(1.82)

        add_rect(sl, left, top + Inches(0.14), Inches(0.07), Inches(0.07), RED)
        f = tb(sl, left + Inches(0.2), top, col_w - Inches(0.2), Inches(1.7))
        p0(f, heading, size=13, bold=True, color=NEAR_BLACK, after=4)
        pn(f, body, size=12, color=DARK_GREY)

    add_rect(sl, ML + col_w + gap/2, Inches(1.45), Inches(0.02), Inches(5.5), GREY_300)

    f = tb(sl, ML, SH - Inches(0.42), CW, Inches(0.32))
    p0(f, "[ Insert: side-by-side Study vs Explore page — annotated with ownership labels ]",
       size=9, color=MID_GREY, italic=True)


def slide_three_cohorts(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Three Cohorts, One Problem")

    cohorts = [
        ("School Leavers", AQUA,
         "Comparing 3–5 institutions under real time pressure.",
         ["Want instant answers — ATAR cutoffs, first-year experience, scholarship eligibility.",
          "When the site doesn't answer immediately, they move on.",
          "The next university that answers them wins."]),
        ("Mature-Age Students", RED,
         "Returning after a career or life break.",
         ["Carry specific anxieties about belonging, eligibility, and whether prior experience counts.",
          "The site treats them the same as school leavers.",
          "No signal that Murdoch understands their situation."]),
        ("Growth Learners", PURPLE,
         "International students & career changers.",
         ["Navigating complex pathway info — credit recognition, English requirements, bridging options.",
          "Content scattered across disconnected pages.",
          "High cognitive load, high-stakes decision."]),
    ]

    cw = Inches(3.88)
    ch = Inches(5.45)
    gap = Inches(0.18)
    ct  = Inches(1.42)

    for i, (name, color, tagline, bullets) in enumerate(cohorts):
        left = ML + i * (cw + gap)
        card(sl, left, ct, cw, ch, LIGHT_GREY, color)

        f = tb(sl, left + Inches(0.25), ct + Inches(0.2), cw - Inches(0.5), Inches(0.5))
        p0(f, name, size=16, bold=True, color=color)

        f = tb(sl, left + Inches(0.25), ct + Inches(0.72), cw - Inches(0.5), Inches(0.55))
        p0(f, tagline, size=12, color=DARK_GREY, italic=True)

        add_rect(sl, left + Inches(0.25), ct + Inches(1.32), cw - Inches(0.5), Inches(0.02), GREY_300)

        f = tb(sl, left + Inches(0.25), ct + Inches(1.45), cw - Inches(0.5), Inches(3.8))
        p0(f, f"• {bullets[0]}", size=12, color=NEAR_BLACK, after=8)
        for b in bullets[1:]:
            pn(f, f"• {b}", size=12, color=NEAR_BLACK, after=8)

    f = tb(sl, ML, SH - Inches(0.42), CW, Inches(0.32))
    p0(f, "[ Insert: journey map showing drop-off points between awareness and enquiry stages ]",
       size=9, color=MID_GREY, italic=True)


def slide_research_stats(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "This Cohort Has Already Moved On From Your Website")

    hw = Inches(5.95)
    ht = Inches(2.9)
    ht_top = Inches(1.45)

    # Stat 1
    card(sl, ML, ht_top, hw, ht, LIGHT_GREY, RED)
    f = tb(sl, ML + Inches(0.3), ht_top + Inches(0.2), hw - Inches(0.5), Inches(1.1))
    p0(f, "2 in 3", size=64, bold=True, color=RED)

    f = tb(sl, ML + Inches(0.3), ht_top + Inches(1.35), hw - Inches(0.5), Inches(0.85))
    p0(f, "prospective students prefer AI-powered search over Google when researching universities", size=13, color=NEAR_BLACK, after=4)
    pn(f, "Everspring, 2025  —  450,000+ student search sessions analysed", size=10, color=MID_GREY, italic=True)

    # Stat 2
    s2l = ML + hw + Inches(0.35)
    card(sl, s2l, ht_top, hw, ht, LIGHT_GREY, AQUA)
    f = tb(sl, s2l + Inches(0.3), ht_top + Inches(0.2), hw - Inches(0.5), Inches(0.95))
    p0(f, "4%  →  23%", size=48, bold=True, color=AQUA)

    f = tb(sl, s2l + Inches(0.3), ht_top + Inches(1.2), hw - Inches(0.5), Inches(1.0))
    p0(f, "high school seniors using AI to research universities (2023 → 2025)", size=13, color=NEAR_BLACK, after=4)
    pn(f, "Carnegie College Choice Trends Report, 2025  —  six-fold growth in two years", size=10, color=MID_GREY, italic=True)

    # Context
    f = tb(sl, ML, Inches(4.55), CW, Inches(1.75))
    p0(f, "Education is the second-highest category for traffic referred by AI tools — behind only online services.", size=15, bold=True, color=NEAR_BLACK, after=8)
    pn(f, "This cohort is not waiting for universities to catch up. They are already asking AI what Murdoch offers, how to get in, and how it compares. The question is whether Murdoch's answers are the ones being surfaced.", size=13, color=DARK_GREY)


def slide_what_we_built(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "What We've Built")

    lw = Inches(5.6)
    f = tb(sl, ML, Inches(1.42), lw, Inches(0.62))
    p0(f, "A working AI agent — built on Claude Sonnet, trained on Murdoch-specific content — that answers the questions prospective students are actually asking.", size=12, color=DARK_GREY)

    bullets = [
        "Streams responses in real time across courses, pathways, fees, entry requirements, and campus life",
        "Adapts to the apparent student cohort — mature-age vs school leaver responses differ accordingly",
        "Maintains context across a conversation — students don't repeat themselves",
        "Surfaces clear next steps: application links, campus tour bookings, advisor handoff",
        "Deployed within the Murdoch brand system — looks and sounds like Murdoch",
    ]

    f = tb(sl, ML, Inches(2.12), lw, Inches(4.55))
    p0(f, f"• {bullets[0]}", size=12, color=NEAR_BLACK, after=10)
    for b in bullets[1:]:
        pn(f, f"• {b}", size=12, color=NEAR_BLACK, after=10)

    screenie(sl, ML + lw + Inches(0.4), Inches(1.42),
             CW - lw - Inches(0.4), Inches(5.42),
             "Murdoch AI agent prototype\nchat interface screenshot")


def slide_live_demo(prs):
    sl = blank(prs)
    fill_bg(sl, DARK_BG)

    f = tb(sl, Inches(0.5), Inches(0.3), SW - Inches(1), Inches(5.5))
    p0(f, "DEMO", size=210, bold=True,
       color=RGBColor(0x2C, 0x2C, 0x2C), align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(4.2), SH/2 - Inches(0.03), Inches(4.9), Inches(0.05), RED)

    f = tb(sl, ML, SH/2 + Inches(0.15), CW, Inches(0.6))
    p0(f, "Live Demonstration", size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    f = tb(sl, Inches(2.2), SH/2 + Inches(0.9), Inches(9.0), Inches(0.55))
    p0(f, "Walking through a live session — showing how the agent handles\nthe queries your prospective students are actually asking.",
       size=14, color=MID_GREY, align=PP_ALIGN.CENTER)


def slide_knowledge_platform(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "A Knowledge Platform, Not a Chatbot")

    f = tb(sl, ML, Inches(1.4), Inches(8.8), Inches(0.42))
    p0(f, "The prototype draws from a curated knowledge base. The production platform connects to authoritative source systems — live data, governed access, no manual updates required.", size=12, color=DARK_GREY)

    rows_data = [
        ("Course Management System",       "Live catalogue, units, prerequisites, fees"),
        ("Admissions / Student Management","Entry requirements, ATAR cutoffs, application status"),
        ("Scholarship Database",           "Eligibility criteria, deadlines, award amounts"),
        ("RPL / Credit Transfer",          "Prior learning rules, credit mapping"),
        ("International Office Systems",   "Visa pathways, English requirements, CRICOS, agent referrals"),
        ("Events & Campus Services",       "Open days, tours, orientation, accommodation"),
        ("CRM (Salesforce / Dynamics)",    "Interaction history, lead status, advisor assignment"),
    ]

    rh   = Inches(0.57)
    ttop = Inches(1.9)
    tw   = Inches(8.7)
    c1w  = Inches(3.8)
    c2w  = tw - c1w

    add_rect(sl, ML, ttop, c1w, rh, NEAR_BLACK)
    add_rect(sl, ML + c1w, ttop, c2w, rh, NEAR_BLACK)

    f = tb(sl, ML + Inches(0.12), ttop + Inches(0.1), c1w, rh)
    p0(f, "Source System", size=12, bold=True, color=WHITE)
    f = tb(sl, ML + c1w + Inches(0.12), ttop + Inches(0.1), c2w, rh)
    p0(f, "What It Enables", size=12, bold=True, color=WHITE)

    for i, (src, enables) in enumerate(rows_data):
        t = ttop + (i + 1) * rh
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(sl, ML, t, c1w, rh, bg)
        add_rect(sl, ML + c1w, t, c2w, rh, bg)
        add_rect(sl, ML, t, tw, Inches(0.01), GREY_300)

        f = tb(sl, ML + Inches(0.12), t + Inches(0.1), c1w - Inches(0.2), rh - Inches(0.1))
        p0(f, src, size=11, bold=True, color=NEAR_BLACK)

        f = tb(sl, ML + c1w + Inches(0.12), t + Inches(0.1), c2w - Inches(0.2), rh - Inches(0.1))
        p0(f, enables, size=11, color=DARK_GREY)

    screenie(sl, ML + tw + Inches(0.35), Inches(1.4),
             CW - tw - Inches(0.35), Inches(5.5),
             "Architecture diagram:\nData sources →\nKnowledge Platform →\nAgent")


def slide_production_capabilities(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Production-Grade From Day One")

    cols_data = [
        ("Conversation Design", AQUA,
         ["Formalised cohort flows", "Escalation triggers",
          "Brand voice governance", "WCAG 2.1 AA accessibility"]),
        ("Integration Layer", RED,
         ["CRM write-back for lead capture", "Application system connection",
          "Live advisor handoff with full context", "Booking integration (tours, appointments)"]),
        ("Quality & Governance", PURPLE,
         ["Grounded responses with source citations", "Hallucination detection & confidence thresholds",
          "Privacy Act compliance & PII handling", "Audit trail & sensitive topic escalation"]),
        ("Analytics Engine", ORANGE,
         ["Real-time dashboard: volume, topics, satisfaction", "Content gap detection from unanswered questions",
          "Funnel attribution: chat → application → enrolment", "A/B testing for conversation flows"]),
    ]

    cw  = Inches(2.93)
    ch  = Inches(5.45)
    gap = Inches(0.14)
    ct  = Inches(1.42)

    for i, (title, color, points) in enumerate(cols_data):
        left = ML + i * (cw + gap)
        card(sl, left, ct, cw, ch, LIGHT_GREY, color)

        f = tb(sl, left + Inches(0.2), ct + Inches(0.18), cw - Inches(0.4), Inches(0.5))
        p0(f, title, size=14, bold=True, color=color)

        add_rect(sl, left + Inches(0.2), ct + Inches(0.76), cw - Inches(0.4), Inches(0.02), GREY_300)

        f = tb(sl, left + Inches(0.2), ct + Inches(0.9), cw - Inches(0.4), ch - Inches(1.1))
        p0(f, f"• {points[0]}", size=11.5, color=NEAR_BLACK, after=9)
        for pt in points[1:]:
            pn(f, f"• {pt}", size=11.5, color=NEAR_BLACK, after=9)


def slide_strategic_case(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "One Capability, Two Strategic Lenses")

    cw  = Inches(5.75)
    ch  = Inches(5.45)
    ct  = Inches(1.42)
    c2l = ML + cw + Inches(0.45)

    # Left — Marketing
    card(sl, ML, ct, cw, ch, LIGHT_GREY, RED)
    f = tb(sl, ML + Inches(0.25), ct + Inches(0.18), cw - Inches(0.5), Inches(0.5))
    p0(f, "The Marketing & Conversion Case", size=14, bold=True, color=RED)
    add_rect(sl, ML + Inches(0.25), ct + Inches(0.76), cw - Inches(0.5), Inches(0.02), GREY_300)

    left_paras = [
        "Murdoch's marketing investment drives traffic to a website that loses prospects before they convert. Every dollar spent on brand, search, and social is partially wasted at the point of information-seeking.",
        "The AI agent is a conversion layer — it captures intent at peak engagement and converts it into action.",
        "When a prospective student gets a clear, confident answer at 10pm on a Sunday while filling out a comparison spreadsheet, that matters. The competitor who answers first wins the consideration set.",
        "The agent also makes visible a feedback loop that doesn't currently exist: every question asked is a signal about what students don't understand, where the narrative breaks down, what content is missing.",
    ]
    f = tb(sl, ML + Inches(0.25), ct + Inches(0.9), cw - Inches(0.5), ch - Inches(1.1))
    p0(f, left_paras[0], size=12, color=NEAR_BLACK, after=8)
    for para_text in left_paras[1:]:
        pn(f, para_text, size=12, color=NEAR_BLACK, after=8)

    # Right — Architecture
    card(sl, c2l, ct, cw, ch, LIGHT_GREY, PURPLE)
    f = tb(sl, c2l + Inches(0.25), ct + Inches(0.18), cw - Inches(0.5), Inches(0.5))
    p0(f, "The Architecture & Platform Case", size=14, bold=True, color=PURPLE)
    add_rect(sl, c2l + Inches(0.25), ct + Inches(0.76), cw - Inches(0.5), Inches(0.02), GREY_300)

    arch_points = [
        "Model-agnostic API layer — updated as capabilities improve without rebuilding",
        "Knowledge layer separated from inference — content teams update without engineering",
        "Integration-ready: CRM, application systems, advisor queues via standard APIs",
        "Audit trail and conversation logging built in for compliance",
        "Flexible deployment: embedded widget, standalone page, or mobile-first interface",
        "Murdoch-owned architecture — no vendor lock-in",
    ]
    f = tb(sl, c2l + Inches(0.25), ct + Inches(0.9), cw - Inches(0.5), ch - Inches(1.1))
    p0(f, f"• {arch_points[0]}", size=12, color=NEAR_BLACK, after=8)
    for pt in arch_points[1:]:
        pn(f, f"• {pt}", size=12, color=NEAR_BLACK, after=8)


def slide_why_now(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "The Window Is Open — Briefly")

    cards_data = [
        ("Students Already Expect It", AQUA,
         "2 in 3 prospective students prefer AI-powered search over Google when researching universities (Everspring, 2025). High school seniors using AI for college search has grown six-fold in two years. This cohort is already asking AI about Murdoch — the question is whether Murdoch's answers are the ones returned."),
        ("The Technology Is Production-Ready", RED,
         "Twelve months ago, deploying a reliable, hallucination-resistant, brand-safe AI agent at institutional scale was genuinely difficult. The models, infrastructure, and integration patterns have matured. The risk of early adoption has dropped substantially."),
        ("Competitors Are Moving", PURPLE,
         "Leading Australian universities are actively piloting AI student experience tools. The question is not whether this becomes standard — it will. The question is whether Murdoch leads that shift in Western Australia or responds to it."),
        ("The Research Is Already Done", ORANGE,
         "Establishing the case normally requires a discovery phase just to prove the problem exists. That work is complete — UX audit, user interviews, journey mapping. The path from insight to implementation is unusually short."),
    ]

    cw  = Inches(5.88)
    ch  = Inches(2.5)
    gap = Inches(0.18)
    positions = [
        (ML,          Inches(1.45)),
        (ML + cw + gap, Inches(1.45)),
        (ML,          Inches(1.45) + ch + gap),
        (ML + cw + gap, Inches(1.45) + ch + gap),
    ]

    for (left, top), (title, color, body) in zip(positions, cards_data):
        card(sl, left, top, cw, ch, LIGHT_GREY, color)

        f = tb(sl, left + Inches(0.25), top + Inches(0.18), cw - Inches(0.5), Inches(0.42))
        p0(f, title, size=13, bold=True, color=color)

        add_rect(sl, left + Inches(0.25), top + Inches(0.68), cw - Inches(0.5), Inches(0.02), GREY_300)

        f = tb(sl, left + Inches(0.25), top + Inches(0.8), cw - Inches(0.5), ch - Inches(0.98))
        p0(f, body, size=11.5, color=NEAR_BLACK)


def slide_roadmap_timeline(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Phased Engagement Roadmap")

    lanes     = ["Knowledge", "Agent", "Integration", "Operations"]
    lane_cols = [AQUA, RED, PURPLE, ORANGE]
    phases    = [
        ("Phase 1 — Prove",    "Weeks 1–8",   RED),
        ("Phase 2 — Federate", "Weeks 9–20",  AQUA),
        ("Phase 3 — Scale",    "Ongoing",     PURPLE),
    ]
    lane_content = [
        ["Build knowledge base\nfrom Murdoch content",    "Connect admissions,\nscholarships & RPL",    "Live data feeds\n& continuous updates"],
        ["Deploy production\nagent to staging",           "Multi-cohort\nconversation flows",           "Full journey coverage\n& authenticated access"],
        ["Advisor handoff\nprotocol",                     "CRM integration\n& lead capture",            "Booking & application\ntracking integration"],
        ["Analytics dashboard\n& baseline",               "Staff training\n& escalation flows",         "Continuous improvement\n& A/B testing"],
    ]
    milestones = ["Pilot Live", "Full Cohort Coverage", "Continuous Optimisation"]
    light_cols = [RED_LIGHT, AQUA_LIGHT, PURPLE_LIGHT]

    lbl_w   = Inches(1.45)
    tl_left = ML + lbl_w + Inches(0.08)
    tl_w    = CW - lbl_w - Inches(0.08)
    ph_w    = tl_w / 3
    lane_h  = Inches(0.9)
    hdr_h   = Inches(0.55)
    top     = Inches(1.45)

    # Phase headers
    for i, (name, weeks, color) in enumerate(phases):
        left = tl_left + i * ph_w
        add_rect(sl, left, top, ph_w - Inches(0.06), hdr_h, color)
        f = tb(sl, left + Inches(0.12), top + Inches(0.04), ph_w - Inches(0.24), hdr_h - Inches(0.06))
        p0(f, name, size=12, bold=True, color=WHITE, after=2)
        pn(f, weeks, size=10, color=RGBColor(0xDD, 0xDD, 0xDD))

    # Lane rows
    for row, (lane_name, lane_color) in enumerate(zip(lanes, lane_cols)):
        t = top + hdr_h + row * lane_h + Inches(0.04)

        add_rect(sl, ML, t, lbl_w - Inches(0.08), lane_h - Inches(0.04), LIGHT_GREY)
        f = tb(sl, ML + Inches(0.1), t + Inches(0.25), lbl_w - Inches(0.2), Inches(0.4))
        p0(f, lane_name, size=12, bold=True, color=lane_color)

        for col, content in enumerate(lane_content[row]):
            left = tl_left + col * ph_w
            lc   = light_cols[col]
            add_rect(sl, left, t, ph_w - Inches(0.06), lane_h - Inches(0.04), lc)
            f = tb(sl, left + Inches(0.1), t + Inches(0.1), ph_w - Inches(0.2), lane_h - Inches(0.2))
            p0(f, content, size=9.5, color=NEAR_BLACK)

    # Milestone markers
    for col, label in enumerate(milestones):
        x = tl_left + (col + 1) * ph_w - Inches(0.04)
        add_rect(sl, x, top + hdr_h, Inches(0.04), len(lanes) * lane_h, NEAR_BLACK)
        f = tb(sl, x - Inches(0.85), top + hdr_h + len(lanes) * lane_h + Inches(0.05),
               Inches(1.75), Inches(0.35))
        p0(f, f"▲ {label}", size=9, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)


def slide_three_phases(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "From Pilot to Platform")

    phases_data = [
        ("Phase 1", "Prove",    "8 weeks",     RED,
         ["Course discovery and entry requirements — highest-friction section",
          "School leaver cohort focus",
          "Website and course system as knowledge sources"],
         ["Production agent deployed to staging",
          "Knowledge base from existing Murdoch content",
          "Analytics dashboard: volume, topics, escalations",
          "Advisor handoff protocol for complex queries"]),
        ("Phase 2", "Federate", "10–12 weeks", AQUA,
         ["Connect admissions, scholarships, and RPL data sources",
          "Expand to mature-age and international cohorts",
          "CRM integration for lead capture and interaction history"],
         ["Federated knowledge platform (4 source systems)",
          "CRM integration & lead capture pipeline",
          "Multi-cohort conversation coverage",
          "Staff training & escalation workflows"]),
        ("Phase 3", "Scale",    "Ongoing",     PURPLE,
         ["Full prospective journey coverage",
          "Authenticated access for returning applicants",
          "Booking integrations & multi-language scoping"],
         ["Full production platform",
          "Application tracking integration",
          "Campus booking integration",
          "Optimisation loop from live conversation data"]),
    ]

    cw  = Inches(3.88)
    ch  = Inches(5.45)
    gap = Inches(0.18)
    ct  = Inches(1.42)

    for i, (phase_num, phase_name, duration, color, focus, deliverables) in enumerate(phases_data):
        left = ML + i * (cw + gap)
        card(sl, left, ct, cw, ch, LIGHT_GREY, color)

        f = tb(sl, left + Inches(0.2), ct + Inches(0.15), cw - Inches(0.4), Inches(0.35))
        p0(f, phase_num, size=11, bold=True, color=color)

        f = tb(sl, left + Inches(0.2), ct + Inches(0.48), cw - Inches(0.4), Inches(0.5))
        p0(f, phase_name, size=24, bold=True, color=NEAR_BLACK)

        add_rect(sl, left + Inches(0.2), ct + Inches(0.98), Inches(1.2), Inches(0.3), color)
        f = tb(sl, left + Inches(0.2), ct + Inches(0.98), Inches(1.2), Inches(0.3))
        p0(f, duration, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_rect(sl, left + Inches(0.2), ct + Inches(1.38), cw - Inches(0.4), Inches(0.02), GREY_300)

        f = tb(sl, left + Inches(0.2), ct + Inches(1.5), cw - Inches(0.4), ch - Inches(2.8))
        p0(f, f"• {focus[0]}", size=11, color=DARK_GREY, after=5)
        for item in focus[1:]:
            pn(f, f"• {item}", size=11, color=DARK_GREY, after=5)

        f = tb(sl, left + Inches(0.2), ct + Inches(2.9), cw - Inches(0.4), Inches(0.32))
        p0(f, "Deliverables", size=11, bold=True, color=color)

        f = tb(sl, left + Inches(0.2), ct + Inches(3.25), cw - Inches(0.4), ch - Inches(3.42))
        p0(f, f"• {deliverables[0]}", size=11, color=NEAR_BLACK, after=5)
        for item in deliverables[1:]:
            pn(f, f"• {item}", size=11, color=NEAR_BLACK, after=5)


def slide_commitment(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "A Clear Commitment to a Defined Scope")

    cw  = Inches(5.75)
    ch  = Inches(5.45)
    ct  = Inches(1.42)
    c2l = ML + cw + Inches(0.45)

    # Left — What We Bring
    card(sl, ML, ct, cw, ch, LIGHT_GREY, RED)
    f = tb(sl, ML + Inches(0.25), ct + Inches(0.18), cw - Inches(0.5), Inches(0.5))
    p0(f, "What We Bring", size=15, bold=True, color=RED)
    add_rect(sl, ML + Inches(0.25), ct + Inches(0.76), cw - Inches(0.5), Inches(0.02), GREY_300)

    bring = [
        "AI product expertise and a track record of deployment at institutional scale",
        "Established integration patterns for university systems (admissions, CRM, LMS)",
        "A working prototype — not a proposal, a product ready to demonstrate",
        "Research evidence that makes this evidence-led from day one",
        "An architecture that Murdoch owns and controls",
    ]
    f = tb(sl, ML + Inches(0.25), ct + Inches(0.9), cw - Inches(0.5), ch - Inches(1.1))
    p0(f, f"• {bring[0]}", size=12, color=NEAR_BLACK, after=10)
    for item in bring[1:]:
        pn(f, f"• {item}", size=12, color=NEAR_BLACK, after=10)

    # Right — What We're Asking For
    card(sl, c2l, ct, cw, ch, NEAR_BLACK)
    add_rect(sl, c2l, ct, cw, Inches(0.06), RED)

    f = tb(sl, c2l + Inches(0.25), ct + Inches(0.18), cw - Inches(0.5), Inches(0.5))
    p0(f, "What We're Asking For", size=15, bold=True, color=RED)
    add_rect(sl, c2l + Inches(0.25), ct + Inches(0.76), cw - Inches(0.5), Inches(0.02), DARK_GREY)

    ask = [
        ("A commitment to Phase 1", True),
        ("Eight weeks", False),
        ("Defined scope: course discovery & entry requirements, school leaver cohort", False),
        ("Shared success criteria: engagement rates, conversion indicators, advisor deflection", False),
        ("Access to course content and admissions data for knowledge base construction", False),
    ]
    f = tb(sl, c2l + Inches(0.25), ct + Inches(0.9), cw - Inches(0.5), ch - Inches(2.0))
    first = True
    for text, bold in ask:
        prefix = "" if bold else "• "
        if first:
            p0(f, f"{prefix}{text}", size=12, bold=bold, color=WHITE, after=10)
            first = False
        else:
            pn(f, f"{prefix}{text}", size=12, bold=bold, color=WHITE, after=10)

    f = tb(sl, c2l + Inches(0.25), ct + ch - Inches(1.05), cw - Inches(0.5), Inches(0.85))
    p0(f, "Eight weeks to prove the concept works — with Murdoch's students, on Murdoch's website, with Murdoch's data.", size=11, color=MID_GREY, italic=True)


def slide_quote(prs):
    sl = blank(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, SW, Inches(0.1), RED)
    add_rect(sl, 0, SH - Inches(0.1), SW, Inches(0.1), RED)

    f = tb(sl, ML, Inches(0.9), Inches(1.5), Inches(1.8))
    p0(f, "\u201c", size=110, bold=True, color=RED)

    f = tb(sl, ML + Inches(1.1), Inches(1.55), Inches(10.8), Inches(2.8))
    p0(f, "Your students are arriving curious and leaving confused. We've built an AI agent that turns that friction into a conversation — and we're ready to prove it works.",
       size=26, color=WHITE)

    add_rect(sl, ML, Inches(4.95), Inches(3.2), Inches(0.05), RED)

    f = tb(sl, ML, Inches(5.2), CW, Inches(0.7))
    p0(f, "The question isn't whether this becomes standard. It's whether Murdoch leads.",
       size=19, bold=True, color=MID_GREY)


def slide_cta(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    add_rect(sl, 0, 0, SW, Inches(0.45), RED)

    f = tb(sl, ML, Inches(1.6), CW, Inches(1.2))
    p0(f, "Let's Talk", size=56, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(5.0), Inches(3.0), Inches(3.3), Inches(0.05), RED)

    f = tb(sl, ML, Inches(3.2), CW, Inches(0.7))
    p0(f, "We're ready to walk through Phase 1 scoping, answer any technical questions, and agree on next steps today.",
       size=15, color=DARK_GREY, align=PP_ALIGN.CENTER)

    f = tb(sl, ML, Inches(4.1), CW, Inches(0.6))
    p0(f, "What questions do you have?", size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

    screenie(sl, Inches(4.7), Inches(5.0), Inches(3.9), Inches(1.9),
             "Optional: closing AI agent screenshot")


def slide_appendix_a(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Research Evidence Summary")

    f = tb(sl, ML, Inches(1.42), CW, Inches(0.28))
    p0(f, "APPENDIX A", size=10, bold=True, color=MID_GREY)

    rows_data = [
        ("UX Audit",            "Navigation maps to internal org structure, not student mental models. Core information consistently difficult to locate."),
        ("UX Audit",            "Study sub-site (Marketing) and Explore sub-site (Schools) operate as separate environments — inconsistent layouts, terminology, depth. Students hit a jarring drop in coherence when crossing between them."),
        ("UX Audit",            "Mobile experience falls significantly below expectation for a primarily mobile-first cohort."),
        ("School Leaver Research", "Students compare 3–5 institutions simultaneously. Whoever answers their question first, clearly, wins."),
        ("School Leaver Research", "Interviewees expressed clear preference for immediate, conversational answers. One participant named AI as their preferred channel over live chat and email — consistent with broader industry research."),
        ("Everspring (2025)",   "2 in 3 prospective students prefer AI-powered search over Google. Education is 2nd highest category for AI-referred traffic. (450,000+ search sessions)"),
        ("Carnegie (2025)",     "High school seniors using AI to research universities: 4% (2023) → 23% (2025). Six-fold growth in two years."),
        ("Mature-Age Research", "Site does not signal that Murdoch understands returning students' specific circumstances or anxieties."),
        ("Growth Learner Research", "Pathway information is fragmented across disconnected pages. High cognitive load, high-stakes decision."),
        ("Journey Map",         "Multiple drop-off points between awareness and enquiry stages — each a recoverable conversion currently being lost."),
    ]

    rh   = Inches(0.48)
    ttop = Inches(1.78)
    c1w  = Inches(2.35)
    c2w  = CW - c1w

    add_rect(sl, ML, ttop, c1w, rh, NEAR_BLACK)
    add_rect(sl, ML + c1w, ttop, c2w, rh, NEAR_BLACK)
    f = tb(sl, ML + Inches(0.1), ttop + Inches(0.1), c1w - Inches(0.2), rh)
    p0(f, "Source", size=11, bold=True, color=WHITE)
    f = tb(sl, ML + c1w + Inches(0.1), ttop + Inches(0.1), c2w - Inches(0.2), rh)
    p0(f, "Key Finding", size=11, bold=True, color=WHITE)

    for i, (src, finding) in enumerate(rows_data):
        t  = ttop + (i + 1) * rh
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(sl, ML, t, c1w, rh, bg)
        add_rect(sl, ML + c1w, t, c2w, rh, bg)
        add_rect(sl, ML, t, CW, Inches(0.01), GREY_300)

        f = tb(sl, ML + Inches(0.1), t + Inches(0.06), c1w - Inches(0.2), rh - Inches(0.08))
        p0(f, src, size=10, bold=True, color=NEAR_BLACK)
        f = tb(sl, ML + c1w + Inches(0.1), t + Inches(0.06), c2w - Inches(0.2), rh - Inches(0.08))
        p0(f, finding, size=10, color=DARK_GREY)


def slide_appendix_b(prs):
    sl = blank(prs)
    fill_bg(sl, WHITE)
    title_bar(sl, "Architecture Overview")

    f = tb(sl, ML, Inches(1.42), CW, Inches(0.28))
    p0(f, "APPENDIX B", size=10, bold=True, color=MID_GREY)

    layers = [
        ("Data Sources", AQUA,
         ["Course Management", "Admissions / SMS",
          "Scholarship Database", "RPL / Credit Transfer",
          "International Office", "Events & Campus Services",
          "CRM (Salesforce / Dynamics)"]),
        ("Knowledge Platform", RED,
         ["Content ingestion & chunking",
          "Semantic embedding & retrieval",
          "Governance & confidence thresholds",
          "Audit trail & logging",
          "Content team updates\n(no engineering required)"]),
        ("Agent Interface", PURPLE,
         ["Conversation engine",
          "Cohort detection & adaptation",
          "Response streaming",
          "Escalation logic",
          "Analytics & reporting"]),
    ]

    bw  = Inches(3.6)
    bh  = Inches(4.35)
    gap = Inches(0.5)
    bt  = Inches(1.85)

    for i, (name, color, items) in enumerate(layers):
        left = ML + i * (bw + gap)
        add_rect(sl, left, bt, bw, bh, LIGHT_GREY, color)
        add_rect(sl, left, bt, bw, Inches(0.5), color)

        f = tb(sl, left + Inches(0.15), bt + Inches(0.08), bw - Inches(0.3), Inches(0.38))
        p0(f, name, size=13, bold=True, color=WHITE)

        f = tb(sl, left + Inches(0.15), bt + Inches(0.62), bw - Inches(0.3), bh - Inches(0.78))
        p0(f, f"• {items[0]}", size=12, color=NEAR_BLACK, after=6)
        for item in items[1:]:
            pn(f, f"• {item}", size=12, color=NEAR_BLACK, after=6)

        if i < len(layers) - 1:
            f = tb(sl, left + bw + Inches(0.1), bt + bh/2 - Inches(0.2),
                   Inches(0.35), Inches(0.4))
            p0(f, "→", size=24, bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Integration bar
    bar_l = ML
    bar_t = bt + bh + Inches(0.22)
    bar_w = 3 * bw + 2 * gap
    bar_h = Inches(0.62)
    add_rect(sl, bar_l, bar_t, bar_w, bar_h, NEAR_BLACK)
    f = tb(sl, bar_l + Inches(0.2), bar_t + Inches(0.1), bar_w - Inches(0.4), bar_h - Inches(0.12))
    p0(f, "Integration layer:  CRM write-back  ·  Advisor handoff with context  ·  Booking systems  ·  Application tracking",
       size=12, color=WHITE, align=PP_ALIGN.CENTER)

    f = tb(sl, bar_l, bar_t + bar_h + Inches(0.15), bar_w, Inches(0.45))
    p0(f, "Model-agnostic  ·  Murdoch-owned architecture  ·  Audit trail for compliance  ·  Flexible deployment: widget · standalone · mobile",
       size=10, color=MID_GREY, italic=True, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_cover(prs)
    slide_agenda(prs)

    slide_section(prs, "01", "The Challenge")
    slide_challenge(prs)
    slide_three_cohorts(prs)
    slide_research_stats(prs)

    slide_section(prs, "02", "The Solution")
    slide_what_we_built(prs)
    slide_live_demo(prs)

    slide_section(prs, "03", "From PoC to Production")
    slide_knowledge_platform(prs)
    slide_production_capabilities(prs)

    slide_section(prs, "04", "The Strategic Case")
    slide_strategic_case(prs)

    slide_section(prs, "05", "Why Now")
    slide_why_now(prs)

    slide_section(prs, "06", "The Roadmap")
    slide_roadmap_timeline(prs)
    slide_three_phases(prs)
    slide_commitment(prs)

    slide_quote(prs)
    slide_cta(prs)

    slide_appendix_a(prs)
    slide_appendix_b(prs)

    out = "/Users/bsummers/claude-code-workbench/pitch/murdoch-ai-pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
